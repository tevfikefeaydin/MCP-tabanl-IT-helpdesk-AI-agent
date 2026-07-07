"""
build_pptx.py
SUNUM.md akisini kurumsal gorunumlu bir PowerPoint (.pptx) dosyasina cevirir.
Konusma notlari PowerPoint'in 'Notlar' bolumune yazilir.

Calistir:
    python tools/build_pptx.py
Cikti:
    ai_helpdesk_agent_demo_sunum.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Renk paleti (kurumsal navy + buz mavisi accent) ---
NAVY = RGBColor(0x0F, 0x1B, 0x2D)
NAVY_SOFT = RGBColor(0x1B, 0x2A, 0x41)
ICE = RGBColor(0x2D, 0xD4, 0xBF)      # teal accent
ICE_DK = RGBColor(0x14, 0x8F, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x22, 0x2E)
SLATE = RGBColor(0x5B, 0x6B, 0x7E)
LOW = RGBColor(0x16, 0xA3, 0x4A)
MED = RGBColor(0xF5, 0x9E, 0x0B)
HIGH = RGBColor(0xDC, 0x26, 0x26)
PAPER = RGBColor(0xF4, 0xF6, 0xF9)

W, H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(shp, color)
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, spacing)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (text, size, color, bold, tracking) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(number, kicker, title, bullets, note, demo, risk=None):
    slide = prs.slides.add_slide(BLANK)
    # Arka plan
    add_rect(slide, 0, 0, W, H, PAPER)
    # Sol accent seridi
    add_rect(slide, 0, 0, Inches(0.28), H, ICE)
    # Ust bant (navy) + slayt no
    add_rect(slide, 0, 0, W, Inches(1.5), NAVY)
    add_rect(slide, 0, 0, Inches(0.28), Inches(1.5), ICE)
    add_text(slide, Inches(0.6), Inches(0.22), Inches(2), Inches(1.1),
             [[(f"{number:02d}", 40, ICE, True, 0)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.7), Inches(0.28), Inches(10.8), Inches(1.0),
             [[(kicker.upper(), 12, ICE, True, 0)], [(title, 26, WHITE, True, 0)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=2)

    # Risk rozeti (varsa)
    if risk:
        rc = {"Low": LOW, "Medium": MED, "High": HIGH}[risk]
        badge = add_rect(slide, Inches(11.0), Inches(0.45), Inches(1.9), Inches(0.6), rc)
        add_text(slide, Inches(11.0), Inches(0.45), Inches(1.9), Inches(0.6),
                 [[(f"RISK: {risk.upper()}", 12, WHITE, True, 0)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Maddeler
    bullet_paras = []
    for b in bullets:
        bullet_paras.append([("•  ", 15, ICE_DK, True, 0), (b, 15, INK, False, 0)])
    add_text(slide, Inches(0.7), Inches(1.85), Inches(7.4), Inches(3.9),
             bullet_paras, space_after=10, line_spacing=1.05)

    # Sag panel: Konusma notu + Demo ekrani
    add_rect(slide, Inches(8.3), Inches(1.85), Inches(4.5), Inches(4.9), WHITE)
    add_rect(slide, Inches(8.3), Inches(1.85), Inches(4.5), Inches(0.08), ICE)
    add_text(slide, Inches(8.6), Inches(2.05), Inches(3.9), Inches(0.4),
             [[("KONUSMA NOTU", 11, ICE_DK, True, 0)]])
    add_text(slide, Inches(8.6), Inches(2.45), Inches(3.9), Inches(2.6),
             [[(note, 12, SLATE, False, 0)]], line_spacing=1.08)
    add_rect(slide, Inches(8.6), Inches(5.35), Inches(3.9), Inches(0.02), RGBColor(0xDD, 0xE3, 0xEA))
    add_text(slide, Inches(8.6), Inches(5.5), Inches(3.9), Inches(0.4),
             [[("DEMO EKRANI", 11, ICE_DK, True, 0)]])
    add_text(slide, Inches(8.6), Inches(5.9), Inches(3.9), Inches(0.8),
             [[(demo, 11.5, INK, False, 0)]], line_spacing=1.05)

    set_notes(slide, f"[{title}]\n\nKONUSMA NOTU:\n{note}\n\nDEMO EKRANI:\n{demo}")
    return slide


def add_box(slide, x, y, w, h, fill, text, tsize=13, tcolor=WHITE, bold=True,
            outline=None, subtitle=None, sub_color=None):
    """Yuvarlak kenarli kutu + ortalanmis metin (opsiyonel alt satir)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if outline:
        shp.line.color.rgb = outline
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(tsize); r.font.bold = bold; r.font.color.rgb = tcolor
    r.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(tsize - 3); r2.font.bold = False
        r2.font.color.rgb = sub_color or tcolor; r2.font.name = "Calibri"
    return shp


def add_down_arrow(slide, cx, y, h, color=ICE_DK, w=Inches(0.22)):
    """Asagi yonlu ok (akis baglantisi)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - w / 2, y, w, h)
    _fill(shp, color)
    return shp


def diagram_slide(number, kicker, title, note, demo):
    """Tam genislik mimari akis diyagrami slayti."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, PAPER)
    add_rect(slide, 0, 0, Inches(0.28), H, ICE)
    add_rect(slide, 0, 0, W, Inches(1.5), NAVY)
    add_rect(slide, 0, 0, Inches(0.28), Inches(1.5), ICE)
    add_text(slide, Inches(0.6), Inches(0.22), Inches(2), Inches(1.1),
             [[(f"{number:02d}", 40, ICE, True, 0)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.7), Inches(0.28), Inches(10.8), Inches(1.0),
             [[(kicker.upper(), 12, ICE, True, 0)], [(title, 26, WHITE, True, 0)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=2)

    cx = Inches(6.66)  # slayt yatay ortasi

    # Row 1: Kullanici / Ticket
    add_box(slide, Inches(5.0), Inches(1.72), Inches(3.33), Inches(0.52),
            NAVY_SOFT, "Kullanici / Ticket", tsize=13)
    add_down_arrow(slide, cx, Inches(2.26), Inches(0.24))

    # Row 2: Streamlit UI
    add_box(slide, Inches(2.5), Inches(2.52), Inches(8.33), Inches(0.56),
            ICE_DK, "Streamlit UI  (app.py)", tsize=14)
    add_down_arrow(slide, cx, Inches(3.10), Inches(0.24))

    # Row 3: AI Agent + iki motor
    add_box(slide, Inches(2.5), Inches(3.36), Inches(8.33), Inches(0.86),
            NAVY, "AI Agent  (ai_agent.py)", tsize=15, subtitle="intent + kategori + tool secimi",
            sub_color=RGBColor(0xAF, 0xC2, 0xD6))
    add_box(slide, Inches(7.05), Inches(3.5), Inches(1.75), Inches(0.28),
            ICE, "Ollama (varsa)", tsize=9.5, tcolor=NAVY)
    add_box(slide, Inches(8.95), Inches(3.5), Inches(1.75), Inches(0.28),
            RGBColor(0xC9, 0xD6, 0xE4), "Rule-based fallback", tsize=9.5, tcolor=NAVY)

    # Uc kola dagilan oklar
    for ax in [Inches(2.55), Inches(6.66), Inches(10.75)]:
        add_down_arrow(slide, ax, Inches(4.24), Inches(0.24))

    # Row 4: Risk Engine | MCP Tools | Response+Audit
    add_box(slide, Inches(0.7), Inches(4.5), Inches(3.7), Inches(1.0),
            NAVY_SOFT, "Risk Engine", tsize=13, outline=HIGH,
            subtitle="Low / Medium / High  ·  High=blocked, Medium=onay",
            sub_color=RGBColor(0xF2, 0xC9, 0xC9))
    add_box(slide, Inches(4.81), Inches(4.5), Inches(3.7), Inches(1.0),
            NAVY_SOFT, "MCP Tools  (11 simule)", tsize=13, outline=ICE,
            subtitle="AI yalnizca tanimli tool'lari cagirir",
            sub_color=RGBColor(0xB9, 0xEE, 0xE7))
    add_box(slide, Inches(8.92), Inches(4.5), Inches(3.7), Inches(1.0),
            NAVY_SOFT, "Response Gen + Audit Logger", tsize=12.5,
            subtitle="Turkce cevap + JSONL log",
            sub_color=RGBColor(0xAF, 0xC2, 0xD6))

    # Row 4 -> Row 5 oklar (iki veri kutusu)
    add_down_arrow(slide, Inches(3.8), Inches(5.54), Inches(0.22))
    add_down_arrow(slide, Inches(9.5), Inches(5.54), Inches(0.22))

    # Row 5: Veri katmani
    add_box(slide, Inches(1.3), Inches(5.78), Inches(5.0), Inches(0.62),
            ICE_DK, "data/*.json", tsize=12.5,
            subtitle="ticket · user · printer · ERP · software · KB",
            sub_color=RGBColor(0xEA, 0xF6, 0xF4))
    add_box(slide, Inches(7.0), Inches(5.78), Inches(5.0), Inches(0.62),
            ICE_DK, "demo_environment/*  (fake)", tsize=12.5,
            subtitle="fake mailbox · desktop · logs",
            sub_color=RGBColor(0xEA, 0xF6, 0xF4))

    # Alt not seridi
    add_text(slide, Inches(0.7), Inches(6.62), Inches(11.9), Inches(0.7),
             [[("Guvenlik kapisi: ", 11, ICE_DK, True, 0),
               ("AI dogrudan sistemlere erisemez; her istek Risk Engine'den gecer, "
                "yalnizca guvenli tool'lar calisir ve her islem loglanir.", 11, SLATE, False, 0)]],
             line_spacing=1.05)

    set_notes(slide, f"[{title}]\n\nKONUSMA NOTU:\n{note}\n\nDEMO EKRANI:\n{demo}")
    return slide


# ---------------------------------------------------------------------------
# Slayt 1 - Kapak
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.18), ICE)
add_rect(s, 0, Inches(7.32), W, Inches(0.18), ICE)
add_text(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.5),
         [[("STAJ SUNUMU  //  KLIMASAN  //  2026", 14, ICE, True, 0)]])
add_text(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.2),
         [[("MCP-Based AI Helpdesk", 46, WHITE, True, 0)],
          [("Agent Demo", 46, WHITE, True, 0)]], space_after=2, line_spacing=1.0)
add_text(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.8),
         [[("Yapay Zeka Destekli, Guvenli ve Local IT Ticket Automation", 20, ICE, False, 0)]])
add_text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
         [[("Gercek sirket verisi kullanilmadi  •  Tamamen local & guvenli  •  Konsept prototipi",
            14, RGBColor(0xAF, 0xC2, 0xD6), False, 0)],
          [("Hazirlayan: Tevfik Efe Aydin", 14, WHITE, True, 0)]], space_after=6)
set_notes(s, "Acilis. Projenin gercek sistemlere baglanmayan, local ve guvenli bir "
             "konsept prototipi oldugunu vurgula.")

# ---------------------------------------------------------------------------
# Icerik slaytlari
# ---------------------------------------------------------------------------
content_slide(
    2, "Problem", "Tekrar Eden Ticket Yuku",
    ["Ticketlarin buyuk kismi standart ve tekrar eden taleplerdir",
     "Ornek: mail eki tasima, yazici sorunu, VPN, disk dolu, program kontrolu",
     "Her ticket: oku - anla - kontrol et - islemi yap - cevap yaz - kapat",
     "Bu dongu zaman alir ve otomasyona cok uygundur"],
    "Bilgi islem personelinin gununun onemli bir kismi bu tur basit ama tekrar "
    "eden ticketlarla geciyor. Sorumuz: bu surecin ne kadarini yapay zeka guvenli "
    "sekilde ustlenebilir?",
    "Uygulamayi ac -> sol paneldeki Ticket dropdown'i goster (8 ornek ticket).")

content_slide(
    3, "Cozum", "AI Agent + Ticket Automation",
    ["AI Agent ticketi okur ve niyetini (intent) anlar",
     "Kategori ve risk seviyesini belirler",
     "Eksik bilgi varsa clarification sorusu uretir",
     "Uygun tool'u secer, guvenliyse islemi yapar",
     "Cevap taslagi uretir, audit log tutar, ticket'i kapatir"],
    "Bu bir sohbet botu degil, gercek bir IT ticket otomasyon prototipi. Agent "
    "ticketi bastan sona bir helpdesk uzmani gibi ele aliyor ve ancak guvenliyse "
    "aksiyon aliyor.",
    "TCK-001 sec -> sol panelde Analyze / Run / Reply / Close / Reset butonlari.")

content_slide(
    4, "Kavram", "MCP ve Tool Calling Nedir?",
    ["MCP: AI'nin dis araclara standart sozlesme ile eristigi mimari",
     "Bu demoda gercek MCP server yok - MCP mantigi simule edildi",
     "Her tool ayri fonksiyon, standart girdi/cikti sozlesmesi var",
     "AI sinirsiz erisemez; yalnizca tanimli tool'lari cagirir",
     "Ornek: search_mailbox, copy_attachment_to_desktop, check_mock_erp_order"],
    "MCP, yapay zeka ile araclar arasindaki koprudur. Kritik ilke: AI dogrudan "
    "dosya sistemine veya sunuculara erisemez; sadece bizim tanimladigimiz sinirli "
    "ve guvenli tool'lari cagirabilir.",
    "TCK-001 -> Analyze Ticket -> AI Analysis'te 'Selected Tool' alanini goster.")

diagram_slide(
    5, "Mimari", "Sistem Mimarisi — Akis Diyagrami",
    "Mimari modulerdir; her katmanin tek sorumlulugu var. En onemli detay: AI ne "
    "siniflandirirsa siniflandirsin, risk karari her zaman bagimsiz Risk Engine "
    "tarafindan yeniden dogrulanir. Bu diyagram, isteklerin nasil katman katman "
    "guvenlik kapilarindan gectigini gosterir.",
    "AI Analysis bolumunu tam goster: Category, Intent, Risk Level, Confidence, "
    "Reasoning Summary.")

content_slide(
    6, "Ana Vurgu", "Guvenlik: Kontrollu Yetenek",
    ["Gercek sirket verisi kullanilmadi (mail, SAP/ERP, PC yok)",
     "Demo tamamen local ve guvenli - her sey fake klasorlerde",
     "AI sistemlere sinirsiz erisemez, sadece tanimli tool'lari kullanir",
     "Riskli islemler otomatik yapilmaz: High = blocked, Medium = onay",
     "Dosya silme tool'u hic yazilmadi; her islem loglanir (Audit Log)"],
    "Bu slayt projenin kalbi. AI agent'larla ilgili en buyuk endise kontrolsuz "
    "erisim. Biz bunu tasarimla cozduk: yuksek riskli hicbir islem otomatik calismaz, "
    "dosya silme yetenegi hic yok. Mesele guc vermek degil, kontrollu guc vermek.",
    "TCK-008 (admin yetkisi) -> Analyze -> Run Safe Action -> Blocked uyarisini goster.",
    risk="High")

content_slide(
    7, "Demo 1", "Mail Eki -> Masaustu (Basari)",
    ["TCK-001: 'Maildeki rapor dosyasini masaustume tasir misiniz?'",
     "Kategori: Mail Attachment Transfer, risk: Low",
     "Tool: copy_attachment_to_desktop -> rapor.xlsx kopyalanir",
     "Turkce cevap taslagi uretilir, ticket Resolved olur",
     "Islem Audit Log'a yazilir"],
    "Uctan uca basarili bir akis. Dikkat: dosya gercekten fake masaustu klasorune "
    "kopyalaniyor, cevap otomatik uretiliyor, ticket kapaniyor ve hepsi loglaniyor.",
    "TCK-001 -> Analyze -> Run Safe Action -> Generate Reply -> Close Ticket. "
    "Demo Environment Preview'da dosyanin masaustunde belirdigini goster.",
    risk="Low")

content_slide(
    8, "Demo 2", "AI Ne Zaman DURUR?",
    ["TCK-002: coklu ekli dosya -> AI islem yapmaz, soru sorar (clarification)",
     "TCK-004: ERP uretim emri -> read-only okuma (Mock ERP)",
     "TCK-003 / TCK-005: Medium risk -> 'Approval recommended'",
     "AI belirsizlikte aksiyon almaz, once netlestirme ister",
     "Read-only islemler guvenlidir ama yine de loglanir"],
    "Iyi bir agent, ne zaman duracagini bilen agent'tir. TCK-002'de AI tahmin "
    "yurutmuyor, kullaniciya soruyor. TCK-004'te sadece okuyor, hicbir sey "
    "degistirmiyor. Bu guvenli otomasyonun temel ilkesi.",
    "TCK-002 -> Analyze -> Run (clarification cikar). Sonra TCK-004 -> Analyze -> "
    "Run (Mock ERP status: 100234 = Open).",
    risk="Medium")

content_slide(
    9, "Teknik", "Her Ortamda Calisan Demo",
    ["AI motoru: Ollama local model (varsa) - ornek qwen3:8b",
     "Model yoksa rule-based fallback devreye girer -> demo bozulmaz",
     "Kod moduler: her fonksiyonun tek sorumlulugu var",
     "Hatalar try/except ile yakalanir, yollar pathlib ile yonetilir",
     "Tek komutla calisir: streamlit run app.py"],
    "Sunum ortaminda internet ya da GPU olmayabilir. Bu yuzden sistemi iki modlu "
    "tasarladim: Ollama varsa gercek LLM, yoksa kural tabanli mod. Demo her ortamda "
    "ayni sekilde calisiyor.",
    "Sayfa ustundeki AI Motoru rozetini goster. Reset Demo Data ile demoyu sifirla.")

content_slide(
    10, "Gelecek", "Konsepten Gercek Sisteme",
    ["Ticket sistemi: Jira / ServiceNow / OTRS entegrasyonu",
     "Mail: Microsoft Graph API ile gercek attachment okuma",
     "Endpoint Automation: Intune / SCCM ile cihaz islemleri",
     "SAP / ERP: read-only API entegrasyonu",
     "Role-based authorization, human approval, SIEM, KPI dashboard"],
    "Bu bir konsept prototipi; amac fikri ve guvenlik mimarisini gostermek. Ayni "
    "yapi gelecekte gercek ticket sistemi, Graph, Intune/SCCM ve SAP/ERP ile entegre "
    "edilebilir; cunku tool tabanli mimari tam da bunun icin tasarlandi. Tesekkurler.",
    "Kapanis. Soru-cevapta merak edilen ticket'i canli gosterebilirsin.")

out = "ai_helpdesk_agent_demo_sunum.pptx"
prs.save(out)
print("Kaydedildi:", out, "| Slayt sayisi:", len(prs.slides._sldIdLst))
